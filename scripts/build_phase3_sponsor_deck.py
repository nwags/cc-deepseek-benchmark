from pathlib import Path
import csv
import shutil
import subprocess
from PIL import Image

from pptx import Presentation
from pptx.util import Inches, Pt
from pptx.dml.color import RGBColor
from pptx.enum.shapes import MSO_SHAPE
from pptx.enum.text import PP_ALIGN, MSO_ANCHOR

ROOT = Path(".")
OUT_DIR = Path("docs/reports/phase3")
PPTX = OUT_DIR / "PHASE3_SPONSOR_DECK.pptx"
PDF = OUT_DIR / "PHASE3_SPONSOR_DECK.pdf"

ASSET_DIR = Path("figures/phase3/deck_assets")
ARCH = ASSET_DIR / "phase_3_router_benchmark_architecture_diagram.png"
CONTAM = ASSET_DIR / "benchmark_contamination_risks_and_mitigations.png"
COST = ASSET_DIR / "phase_3_canary_cost_breakdown.png"
METRICS = ASSET_DIR / "phase_3_metrics_plan_overview.png"
LEDGER = Path("results/phase3/supplemental/canary_ledger.csv")

for path in [ARCH, CONTAM, COST, METRICS, LEDGER]:
    if not path.exists():
        raise SystemExit(f"missing required input: {path}")

prs = Presentation()
prs.slide_width = Inches(13.333)
prs.slide_height = Inches(7.5)

W, H = 13.333, 7.5
NAVY = "0B1020"
BLUE = "1D4ED8"
CYAN = "38BDF8"
GREEN = "22C55E"
AMBER = "F59E0B"
RED = "EF4444"
WHITE = "FFFFFF"
SLATE = "475569"
LIGHT = "F8FAFC"
MID = "CBD5E1"

def rgb(hexstr):
    return RGBColor.from_string(hexstr)

def blank_slide():
    slide = prs.slides.add_slide(prs.slide_layouts[6])
    bg = slide.background.fill
    bg.solid()
    bg.fore_color.rgb = rgb(LIGHT)
    return slide

def add_text(slide, text, x, y, w, h, size=18, color=NAVY, bold=False, align=None):
    box = slide.shapes.add_textbox(Inches(x), Inches(y), Inches(w), Inches(h))
    tf = box.text_frame
    tf.clear()
    tf.word_wrap = True
    tf.vertical_anchor = MSO_ANCHOR.TOP
    p = tf.paragraphs[0]
    p.text = text
    p.font.size = Pt(size)
    p.font.bold = bold
    p.font.color.rgb = rgb(color)
    p.font.name = "Aptos"
    if align:
        p.alignment = align
    return box

def add_title(slide, title, subtitle=None):
    add_text(slide, title, 0.55, 0.25, 12.3, 0.45, 25, NAVY, True)
    if subtitle:
        add_text(slide, subtitle, 0.58, 0.78, 12.1, 0.35, 12, SLATE)

def add_footer(slide, idx):
    add_text(slide, f"Phase 3 Router Benchmark | {idx}", 0.55, 7.13, 3.0, 0.2, 8, SLATE)
    add_text(slide, "Evidence-backed sponsor briefing", 9.4, 7.13, 3.3, 0.2, 8, SLATE, align=PP_ALIGN.RIGHT)

def card(slide, x, y, w, h, title, value, detail, color=BLUE, link=None):
    sh = slide.shapes.add_shape(MSO_SHAPE.ROUNDED_RECTANGLE, Inches(x), Inches(y), Inches(w), Inches(h))
    sh.fill.solid()
    sh.fill.fore_color.rgb = rgb(WHITE)
    sh.line.color.rgb = rgb(MID)
    if link:
        sh.click_action.hyperlink.address = link
    add_text(slide, title, x + 0.15, y + 0.13, w - 0.3, 0.25, 10, SLATE, True)
    add_text(slide, value, x + 0.15, y + 0.42, w - 0.3, 0.42, 25, color, True)
    add_text(slide, detail, x + 0.15, y + 0.93, w - 0.3, h - 1.05, 10, NAVY)
    return sh

def bullet_list(slide, items, x, y, w, h, size=14):
    box = slide.shapes.add_textbox(Inches(x), Inches(y), Inches(w), Inches(h))
    tf = box.text_frame
    tf.clear()
    tf.word_wrap = True
    for i, item in enumerate(items):
        p = tf.paragraphs[0] if i == 0 else tf.add_paragraph()
        p.text = item
        p.level = 0
        p.font.size = Pt(size)
        p.font.name = "Aptos"
        p.font.color.rgb = rgb(NAVY)
    return box

def image_contain(slide, path, x, y, w, h):
    with Image.open(path) as im:
        iw, ih = im.size
    ir = iw / ih
    br = w / h
    if br > ir:
        nh = h
        nw = h * ir
    else:
        nw = w
        nh = w / ir
    nx = x + (w - nw) / 2
    ny = y + (h - nh) / 2
    return slide.shapes.add_picture(str(path), Inches(nx), Inches(ny), width=Inches(nw), height=Inches(nh))

def link_chip(slide, label, target, x, y, w):
    sh = slide.shapes.add_shape(MSO_SHAPE.ROUNDED_RECTANGLE, Inches(x), Inches(y), Inches(w), Inches(0.32))
    sh.fill.solid()
    sh.fill.fore_color.rgb = rgb("E0F2FE")
    sh.line.color.rgb = rgb(CYAN)
    sh.click_action.hyperlink.address = target
    add_text(slide, label, x + 0.08, y + 0.07, w - 0.16, 0.18, 8, BLUE, True, PP_ALIGN.CENTER)

def rows_from_ledger():
    rows = []
    with LEDGER.open() as f:
        for r in csv.DictReader(f):
            if r.get("classification") == "PASS":
                rows.append(r)
    return rows

# Slide 1
s = blank_slide()
add_title(s, "Phase 3 Router Benchmark: Canary-Green, Multi-Provider, Ready for Funded Smoke",
          "Claude Code remains fixed; LiteLLM routes the same benchmark flow across 13 canary-green arms spanning 8 provider families.")
card(s, 0.65, 1.35, 2.85, 1.45, "READY ARMS", "13", "Active router arms passed the canary task.", GREEN, "docs/reports/phase3/PHASE3_CANARY_EVIDENCE.md")
card(s, 3.75, 1.35, 2.85, 1.45, "PROVIDER FAMILIES", "8", "Anthropic, DeepSeek, Gemini, OpenAI, xAI, Moonshot, DashScope, Z.AI.", BLUE, "configs/arms/")
card(s, 6.85, 1.35, 2.85, 1.45, "5-TASK SMOKE", "$44.72", "Canary-scaled estimate; recommended 2x reserve: $89.44.", AMBER, "results/phase3/supplemental/phase3_cost_forecast_summary.json")
card(s, 9.95, 1.35, 2.85, 1.45, "CONTAMINATION AUDIT", "Clean", "No WebSearch/WebFetch events in passing canaries.", GREEN, "docs/reports/phase3/PHASE3_CANARY_EVIDENCE.md")
add_text(s, "Decision path", 0.75, 3.35, 2.4, 0.3, 16, NAVY, True)
steps = [("1", "Fund 5-task smoke"), ("2", "Replace estimates with smoke data"), ("3", "Approve, resize, or defer full sweep")]
for i, (n, t) in enumerate(steps):
    x = 0.75 + i * 4.15
    card(s, x, 3.8, 3.25, 1.05, f"STEP {n}", t, "", BLUE)
add_text(s, "Core message: Phase 3 is operational enough for funded smoke testing, but the full sweep should wait until smoke results replace canary-scaled estimates.", 0.8, 5.25, 11.8, 0.55, 17, NAVY, True)
link_chip(s, "Evidence ledger", "docs/reports/phase3/PHASE3_CANARY_EVIDENCE.md", 0.8, 6.15, 1.5)
link_chip(s, "Cost summary", "results/phase3/supplemental/phase3_cost_forecast_summary.json", 2.5, 6.15, 1.5)
link_chip(s, "Dashboard spec", "docs/reports/phase3/PHASE3_EXECUTIVE_DASHBOARD_SPEC.md", 4.2, 6.15, 1.7)
add_footer(s, 1)

# Slide 2
s = blank_slide()
add_title(s, "From Three-Arm Baseline to Multi-Provider Router Harness")
card(s, 0.8, 1.35, 5.3, 3.8, "PHASE 1 BASELINE", "3 arms", "Anthropic Sonnet, DeepSeek Pro, DeepSeek Flash. Claude Code fixed as the agent harness. Harbor ran Terminal-Bench tasks.", BLUE)
card(s, 7.2, 1.35, 5.3, 3.8, "PHASE 3 ROUTER LAYER", "13 green arms", "LiteLLM routes the same benchmark flow across 8 provider families. Canary-first execution catches auth, route, quota, schema, and tool-mode issues.", GREEN)
bullet_list(s, ["Continuity: Claude Code remains the fixed harness.", "Change: model/provider coverage expands through LiteLLM.", "Methodology: results remain separated under results/phase3."], 1.0, 5.55, 11.6, 0.9, 13)
add_footer(s, 2)

# Slide 3
s = blank_slide()
add_title(s, "Phase 3 Router Benchmark Architecture", "Request and response flow is bidirectional between Claude Code, LiteLLM, and provider backends.")
image_contain(s, ARCH, 0.55, 1.05, 12.2, 5.8)
add_footer(s, 3)

# Slide 4
s = blank_slide()
add_title(s, "Canary Gate: Infrastructure Readiness, Not Final Model Ranking")
bullet_list(s, [
    "Current gate uses Terminal-Bench 2.0 task modernize-scientific-stack.",
    "The task exercises read, write, execute, verify, routing, results, and audit flow.",
    "A canary pass means routing/config/tool controls are ready for smoke.",
    "It does not rank model quality across the full Terminal-Bench distribution."
], 0.75, 1.35, 5.4, 4.9, 15)
card(s, 7.0, 1.35, 5.3, 1.15, "PASS SIGNAL", "End-to-end path works", "Harbor -> Claude Code -> LiteLLM -> provider -> results -> audit", GREEN)
card(s, 7.0, 2.8, 5.3, 1.15, "FAIL SIGNAL", "Classify before spending", "Auth, quota, route, model slug, schema, tool-mode, or benchmark failure", AMBER)
card(s, 7.0, 4.25, 5.3, 1.15, "NEXT GATE", "5-task smoke", "Use representative multi-task evidence to replace canary-scaled cost estimates", BLUE)
add_footer(s, 4)

# Slide 5
s = blank_slide()
add_title(s, "Current Canary-Green Model Arms", "13 active router arms spanning 8 provider families.")
rows = rows_from_ledger()
x0, y0 = 0.55, 1.15
headers = ["Arm", "Provider", "Backend", "Cost", "Runtime"]
widths = [3.0, 1.7, 3.1, 1.1, 1.1]
for i, h in enumerate(headers):
    add_text(s, h, x0 + sum(widths[:i]), y0, widths[i], 0.25, 9, WHITE, True)
    rect = s.shapes.add_shape(MSO_SHAPE.RECTANGLE, Inches(x0 + sum(widths[:i]) - 0.02), Inches(y0 - 0.05), Inches(widths[i]), Inches(0.35))
    rect.fill.solid(); rect.fill.fore_color.rgb = rgb(BLUE); rect.line.color.rgb = rgb(BLUE)
    rect.z_order = 0
for j, r in enumerate(rows[:13]):
    y = y0 + 0.42 + j * 0.38
    vals = [
        r.get("arm_id", ""),
        r.get("provider", ""),
        r.get("backend_model", ""),
        "$" + str(round(float(r.get("cost_usd", r.get("cost", 0)) or 0), 3)),
        str(round(float(r.get("runtime_seconds", r.get("runtime", 0)) or 0), 1)) + "s",
    ]
    for i, val in enumerate(vals):
        add_text(s, val, x0 + sum(widths[:i]), y, widths[i], 0.25, 8, NAVY)
add_footer(s, 5)

# Slide 6
s = blank_slide()
add_title(s, "Provider-Specific Engineering Lessons")
lessons = [
    ("Gemini", "Route correction and no-plan tool mitigation were needed."),
    ("OpenAI", "Router validation completed; costs need smoke confirmation."),
    ("Qwen", "Singapore endpoint and model-access normalization were required."),
    ("xAI Grok", "Funding/access recovery and model selection unblocked canary."),
    ("Moonshot Kimi", "Temperature/max-token behavior matters for direct probes."),
    ("Z.AI GLM", "Funding/access recovery and model update unblocked canary.")
]
for i, (name, detail) in enumerate(lessons):
    x = 0.7 + (i % 2) * 6.05
    y = 1.25 + (i // 2) * 1.55
    card(s, x, y, 5.45, 1.15, name, "Resolved", detail, GREEN)
add_text(s, "Interpretation: these were infrastructure learnings unless the model actually ran and failed the benchmark task.", 0.8, 6.35, 11.6, 0.35, 14, NAVY, True)
add_footer(s, 6)

# Slide 7
s = blank_slide()
add_title(s, "Benchmark Contamination Controls", "Contamination risk is treated as a first-class control, not an afterthought.")
image_contain(s, CONTAM, 0.55, 1.05, 12.2, 5.8)
add_footer(s, 7)

# Slide 8
s = blank_slide()
add_title(s, "Canary Cost Breakdown", "Canary costs reveal budget drivers before smoke/full-sweep funding.")
image_contain(s, COST, 0.55, 1.05, 12.2, 5.8)
add_footer(s, 8)

# Slide 9
s = blank_slide()
add_title(s, "Funding Gate and Forecast")
card(s, 0.8, 1.3, 3.6, 1.5, "5-TASK SMOKE", "$44.72", "Recommended reserve: $89.44. Run first.", GREEN, "results/phase3/supplemental/phase3_cost_forecast_summary.json")
card(s, 4.85, 1.3, 3.6, 1.5, "20 TASKS x 3", "$536.64", "Reserve: $804.97. Approve after smoke.", AMBER, "results/phase3/supplemental/phase3_cost_forecast_summary.json")
card(s, 8.9, 1.3, 3.6, 1.5, "25 TASKS x 3", "$670.80", "Reserve: $1,006.21. Expanded coverage option.", BLUE, "results/phase3/supplemental/phase3_cost_forecast_summary.json")
bullet_list(s, [
    "Canary-scaled estimates are useful for order-of-magnitude planning only.",
    "Smoke results should become the new source of truth.",
    "GPT-5.5 and Gemini Flash were high-cost canary outliers.",
    "The ask is to fund the next statistically and financially sensible gate."
], 1.0, 3.45, 11.4, 1.6, 16)
add_footer(s, 9)

# Slide 10
s = blank_slide()
add_title(s, "Recommended Next Actions")
items = [
    ("1", "Confirm exact 5-task smoke set"),
    ("2", "Approve smoke budget plus reserve"),
    ("3", "Run smoke across 13 active green router arms"),
    ("4", "Re-run evidence extraction after smoke"),
    ("5", "Replace canary-scaled estimates with smoke-scaled estimates"),
    ("6", "Approve, resize, or defer full sweep")
]
for i, (n, t) in enumerate(items):
    x = 0.9 + (i % 2) * 6.0
    y = 1.25 + (i // 2) * 1.35
    card(s, x, y, 5.4, 0.95, f"ACTION {n}", t, "", BLUE)
add_footer(s, 10)

# Appendix A
s = blank_slide()
add_title(s, "Appendix A: Metrics Collected and Planned")
image_contain(s, METRICS, 0.55, 1.05, 12.2, 5.8)
add_footer(s, 11)

# Appendix B
s = blank_slide()
add_title(s, "Appendix B: Evidence and Reproducibility Hub")
links = [
    "docs/reports/phase3/PHASE3_CANARY_EVIDENCE.md",
    "results/phase3/supplemental/canary_ledger.csv",
    "results/phase3/supplemental/canary_ledger.json",
    "results/phase3/supplemental/phase3_cost_forecast_canary_scaled.csv",
    "results/phase3/supplemental/phase3_cost_forecast_summary.json",
    "configs/router/litellm.config.yaml.example",
    "configs/arms/",
    "BENCHMARK_CONTAMINATION.md",
    "RUNBOOK.md",
    "ARTIFACT_POLICY.md",
]
for i, link in enumerate(links):
    x = 0.9 + (i % 2) * 6.0
    y = 1.25 + (i // 2) * 0.55
    link_chip(s, link, link, x, y, 5.4)
add_footer(s, 12)

OUT_DIR.mkdir(parents=True, exist_ok=True)
prs.save(PPTX)
print(f"wrote {PPTX}")

converter = shutil.which("soffice") or shutil.which("libreoffice")
if not converter:
    raise SystemExit("LibreOffice/soffice not found; PPTX was created but PDF export cannot run.")

subprocess.run([
    converter,
    "--headless",
    "--convert-to", "pdf",
    "--outdir", str(OUT_DIR),
    str(PPTX),
], check=True)
if not PDF.exists():
    raise SystemExit(f"expected PDF not found: {PDF}")
print(f"wrote {PDF}")
